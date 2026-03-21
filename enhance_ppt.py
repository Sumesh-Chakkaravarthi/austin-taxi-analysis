import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def add_title_slide(prs):
    # Use Title Slide layout (typically index 0)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Austin Taxi Analysis - Portfolio Highlights"
    
    # Add a custom textbox for the subtitle
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
    tf = txBox.text_frame
    tf.text = "By Sumesh | GitHub: Sumesh-Chakkaravarthi/austin-taxi-analysis\n\n- Python EDA & Geospatial Visualization\n- Gradient Boosting ML Pipeline (R² = 0.9964)\n- GitHub Actions CI/CD & Automated Reporting"
    for p in tf.paragraphs:
        p.font.size = Pt(20)
    
    # move to index 27 (Right before Slide 28 Austin Dataset Overview)
    move_slide(prs, len(prs.slides) - 1, 27)

def add_visual_slide(prs, title_text, img_paths, insert_idx):
    # Use blank layout or title only layout (typically index 5 or 6)
    slide_layout = prs.slide_layouts[5] # usually Title Only
    slide = prs.slides.add_slide(slide_layout)
    
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = title_text
    else:
        # Add textbox for title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
    
    # Calculate positions for images
    # Slide width typically 10 or 13.3 inches, height 7.5 inches
    prs_width = prs.slide_width
    prs_height = prs.slide_height
    
    num_imgs = len(img_paths)
    if num_imgs == 1:
        # Center the single image
        img_path = img_paths[0]
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(2), Inches(1.5), height=Inches(5.5))
    elif num_imgs == 2:
        img_path1 = img_paths[0]
        img_path2 = img_paths[1]
        
        # side by side
        if os.path.exists(img_path1):
            slide.shapes.add_picture(img_path1, Inches(0.5), Inches(1.5), height=Inches(5))
        if prs_width > Inches(11): # 16:9
            left2 = Inches(6.5)
        else:
            left2 = Inches(5.0)
            
        if os.path.exists(img_path2):
            slide.shapes.add_picture(img_path2, left2, Inches(1.5), height=Inches(5))
            
    # Move slide
    move_slide(prs, len(prs.slides) - 1, insert_idx)

def main():
    prs_path = "Master_Deck_Final_group_presentation 1.pptx"
    out_path = "Master_Deck_Final_group_presentation_Enhanced.pptx"
    
    if not os.path.exists(prs_path):
        print(f"File {prs_path} not found.")
        return
        
    prs = Presentation(prs_path)
    
    # To keep track of our inserts, we need to adjust indexes as we insert
    # Let's insert from back to front to avoid index shifting problems for prior slides
    
    # Slide 46 is index 45. We want to insert after it, so index 46. (We'll do this first since it's the latest)
    # ML Slide
    add_visual_slide(prs, "ML Model Performance & Feature Importance", 
                     ["outputs/12_feature_importance.png", "outputs/14_actual_vs_predicted.png"], 
                     46)
                     
    # Slide 32 is index 31. We want to insert after it, so index 32.
    add_visual_slide(prs, "Surge Pricing & Spatial Demand", 
                     ["outputs/05_surge_analysis.png", "outputs/11_demand_heatmap.png"], 
                     32)
                     
    # Slide 31 is index 30. Insert at index 31.
    add_visual_slide(prs, "Temporal Demand Patterns", 
                     ["outputs/03_hourly_demand.png", "outputs/04_day_of_week.png"], 
                     31)
                     
    # Slide 30 is index 29. Insert at index 30.
    add_visual_slide(prs, "Fare & Distance Distributions", 
                     ["outputs/01_fare_distribution.png", "outputs/02_distance_distribution.png"], 
                     30)
                     
    # Finally, insert the intro slide at index 27
    add_title_slide(prs)
    
    prs.save(out_path)
    print(f"Enhanced presentation saved to {out_path}")

if __name__ == "__main__":
    main()
